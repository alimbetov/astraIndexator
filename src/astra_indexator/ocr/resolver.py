from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Protocol

from docx import Document
from PIL import Image
from pptx import Presentation

from astra_indexator.acquisition import AcquiredSource
from astra_indexator.parser import OcrCandidate, ParsedDocument

from .model import OcrProfile, ResolvedOcrInput


class OcrInputResolver(Protocol):
    def resolve(
        self,
        *,
        source: AcquiredSource,
        document: ParsedDocument,
        candidate: OcrCandidate,
        profile: OcrProfile,
    ) -> ResolvedOcrInput: ...


class DefaultOcrInputResolver:
    def __init__(self, workspace_root: Path):
        self.workspace_root = workspace_root

    def _target(self, document: ParsedDocument, candidate: OcrCandidate) -> Path:
        """Allocate a unique attempt-local scratch path.

        OCR inputs are disposable scratch artifacts. A deterministic candidate name
        alone is unsafe on a shared/reclaimed workspace because two attempts can
        overlap briefly around lease loss. mkstemp gives each resolution a unique
        path while keeping a useful document/candidate prefix for diagnostics.
        """
        target_dir = self.workspace_root / "ocr-inputs"
        target_dir.mkdir(parents=True, exist_ok=True)
        candidate_safe = candidate.candidate_id.replace(":", "_").replace("/", "_")[:80]
        doc_prefix = f"{document.document_id}-{document.document_version}-{candidate_safe}-"
        fd, name = tempfile.mkstemp(prefix=doc_prefix, suffix=".png", dir=target_dir)
        Path(name).unlink(missing_ok=True)
        try:
            import os
            os.close(fd)
        except OSError:
            pass
        return Path(name)

    @staticmethod
    def _dimensions(path: Path) -> tuple[int, int]:
        with Image.open(path) as image:
            return image.width, image.height

    @staticmethod
    def _crop_region(image: Image.Image, candidate: OcrCandidate) -> Image.Image:
        if candidate.scope != "REGION" or candidate.geometry is None:
            return image
        g = candidate.geometry
        if None in (g.x0, g.y0, g.x1, g.y1):
            return image
        if None not in (g.page_width, g.page_height) and g.page_width and g.page_height:
            sx = image.width / float(g.page_width)
            sy = image.height / float(g.page_height)
            box = (int(g.x0 * sx), int(g.y0 * sy), int(g.x1 * sx), int(g.y1 * sy))
        elif g.coordinate_space == "normalized-0-1":
            box = (int(g.x0 * image.width), int(g.y0 * image.height), int(g.x1 * image.width), int(g.y1 * image.height))
        else:
            return image
        left, top, right, bottom = box
        left = max(0, min(left, image.width))
        right = max(0, min(right, image.width))
        top = max(0, min(top, image.height))
        bottom = max(0, min(bottom, image.height))
        if right <= left or bottom <= top:
            raise RuntimeError("OCR_REGION_GEOMETRY_INVALID")
        return image.crop((left, top, right, bottom))

    def resolve(self, *, source, document, candidate, profile) -> ResolvedOcrInput:
        fmt = source.detected_format.upper()
        target = self._target(document, candidate)
        try:
            if fmt in {"JPEG", "PNG", "TIFF"}:
                with Image.open(source.local_path) as image:
                    frame = max(0, (candidate.page_number or 1) - 1)
                    try:
                        image.seek(frame)
                    except EOFError as exc:
                        raise RuntimeError("OCR_IMAGE_FRAME_NOT_FOUND") from exc
                    rendered = self._crop_region(image.convert("RGB"), candidate)
                    rendered.save(target, "PNG")
            elif fmt == "PDF" and candidate.scope in {"PAGE", "REGION"}:
                try:
                    import pypdfium2 as pdfium
                except ImportError as exc:  # pragma: no cover
                    raise RuntimeError("pypdfium2 is required for PDF OCR; install astra-indexator[ocr]") from exc
                page_number = candidate.page_number or 1
                pdf = pdfium.PdfDocument(str(source.local_path))
                try:
                    if page_number < 1 or page_number > len(pdf):
                        raise RuntimeError("OCR_PDF_PAGE_NOT_FOUND")
                    page = pdf[page_number - 1]
                    scale = profile.render_dpi / 72.0
                    bitmap = page.render(scale=scale)
                    image = bitmap.to_pil()
                    image = self._crop_region(image, candidate)
                    image.convert("RGB").save(target, "PNG")
                finally:
                    pdf.close()
            elif fmt == "DOCX" and candidate.scope == "EMBEDDED_IMAGE":
                doc = Document(str(source.local_path))
                paragraph_index = int(candidate.metadata["paragraphIndex"])
                drawing_index = int(candidate.metadata["drawingIndex"])
                if paragraph_index < 0 or paragraph_index >= len(doc.paragraphs):
                    raise RuntimeError("OCR_DOCX_PARAGRAPH_NOT_FOUND")
                paragraph = doc.paragraphs[paragraph_index]
                drawings = paragraph._p.xpath(".//w:drawing")
                if drawing_index < 0 or drawing_index >= len(drawings):
                    raise RuntimeError("OCR_DOCX_DRAWING_NOT_FOUND")
                blips = drawings[drawing_index].xpath(".//a:blip")
                if not blips:
                    raise RuntimeError("OCR_DOCX_IMAGE_RELATION_NOT_FOUND")
                relation_id = blips[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                if not relation_id or relation_id not in doc.part.related_parts:
                    raise RuntimeError("OCR_DOCX_IMAGE_RELATION_NOT_FOUND")
                part = doc.part.related_parts[relation_id]
                with Image.open(__import__("io").BytesIO(part.blob)) as image:
                    image.convert("RGB").save(target, "PNG")
            elif fmt == "PPTX" and candidate.scope == "EMBEDDED_IMAGE":
                presentation = Presentation(str(source.local_path))
                slide_number = int(candidate.metadata["slideNumber"])
                shape_id = int(candidate.metadata["shapeId"])
                if slide_number < 1 or slide_number > len(presentation.slides):
                    raise RuntimeError("OCR_PPTX_SLIDE_NOT_FOUND")
                shape = next((s for s in presentation.slides[slide_number - 1].shapes if int(s.shape_id) == shape_id), None)
                if shape is None or not hasattr(shape, "image"):
                    raise RuntimeError("OCR_PPTX_IMAGE_NOT_FOUND")
                with Image.open(__import__("io").BytesIO(shape.image.blob)) as image:
                    image.convert("RGB").save(target, "PNG")
            else:
                raise RuntimeError(f"OCR_INPUT_UNSUPPORTED:{fmt}:{candidate.scope}")
            width, height = self._dimensions(target)
            return ResolvedOcrInput(
                candidate_id=candidate.candidate_id,
                image_path=target,
                width=width,
                height=height,
                page_number=candidate.page_number,
                source_element_id=candidate.element_id,
                source_geometry=candidate.geometry,
                cleanup=True,
            )
        except Exception:
            target.unlink(missing_ok=True)
            raise
