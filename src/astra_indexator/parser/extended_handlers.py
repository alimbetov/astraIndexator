from __future__ import annotations

import csv
from typing import Any

from bs4 import BeautifulSoup, Tag
from ebooklib import ITEM_DOCUMENT, epub
from odf import teletype
from odf.opendocument import load as load_odf
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from striprtf.striprtf import rtf_to_text

from astra_indexator.acquisition import AcquiredSource

from .base import ParseContext, ParserError
from .handlers import _doc, _eid
from .model import DocumentElement, ElementType, OcrCandidate, SourceGeometry


def _append(
    elements: list[DocumentElement],
    source: AcquiredSource,
    context: ParseContext,
    kind: ElementType,
    locator: str,
    *,
    text: str | None = None,
    metadata: dict[str, Any] | None = None,
    source_locator: dict[str, Any] | None = None,
    geometry: SourceGeometry | None = None,
    role: str | None = None,
    level: int | None = None,
    section_path: tuple[str, ...] = (),
) -> DocumentElement:
    if len(elements) >= context.limits.max_elements:
        raise ParserError("PARSER_ELEMENT_LIMIT", "parser element limit exceeded")
    element = DocumentElement(
        element_id=_eid(source, context, locator, kind),
        type=kind,
        order_index=len(elements),
        text=text,
        level=level,
        geometry=geometry,
        section_path=section_path,
        source_locator=source_locator or {},
        role=role,
        metadata=metadata or {},
    )
    elements.append(element)
    return element


class XlsxDocumentHandler:
    format_name = "XLSX"

    def parse(self, source: AcquiredSource, context: ParseContext):
        formula_wb = load_workbook(
            source.local_path, read_only=True, data_only=False, keep_links=False
        )
        value_wb = load_workbook(
            source.local_path, read_only=True, data_only=True, keep_links=False
        )
        try:
            if len(formula_wb.sheetnames) > context.limits.max_sheets:
                raise ParserError("PARSER_SHEET_LIMIT", "workbook sheet limit exceeded")
            elements: list[DocumentElement] = []
            non_empty_cells = 0
            for sheet_index, sheet_name in enumerate(formula_wb.sheetnames):
                ws_formula = formula_wb[sheet_name]
                ws_value = value_wb[sheet_name]
                if ws_formula.max_column > context.limits.max_columns:
                    raise ParserError("PARSER_COLUMN_LIMIT", "worksheet column limit exceeded")
                headers: list[str] | None = None
                emitted_rows = 0
                for row_index, (formula_row, value_row) in enumerate(
                    zip(ws_formula.iter_rows(), ws_value.iter_rows()), start=1
                ):
                    if row_index > context.limits.max_rows:
                        raise ParserError("PARSER_ROW_LIMIT", "worksheet row limit exceeded")
                    cells: list[dict[str, Any]] = []
                    values_for_text: list[str] = []
                    for f_cell, v_cell in zip(formula_row, value_row):
                        raw = f_cell.value
                        cached = v_cell.value
                        if raw is None and cached is None:
                            continue
                        non_empty_cells += 1
                        if non_empty_cells > context.limits.max_non_empty_cells:
                            raise ParserError(
                                "PARSER_CELL_LIMIT", "workbook non-empty cell limit exceeded"
                            )
                        formula_text = raw if isinstance(raw, str) and raw.startswith("=") else None
                        display_value = (
                            cached if cached is not None else (None if formula_text else raw)
                        )
                        rendered = "" if display_value is None else str(display_value)
                        if len(rendered) > context.limits.max_cell_chars:
                            raise ParserError(
                                "PARSER_CELL_TEXT_LIMIT", "cell text exceeds configured limit"
                            )
                        cells.append(
                            {
                                "address": f_cell.coordinate,
                                "rowIndex": f_cell.row,
                                "columnIndex": f_cell.column,
                                "displayValue": display_value,
                                "formulaPresent": formula_text is not None,
                                "formulaText": formula_text,
                                "numberFormat": f_cell.number_format,
                            }
                        )
                        values_for_text.append(rendered)
                    if not cells:
                        continue
                    emitted_rows += 1
                    if headers is None:
                        headers = values_for_text
                    paired = []
                    if headers and row_index > 1:
                        for idx, value in enumerate(values_for_text):
                            key = (
                                headers[idx]
                                if idx < len(headers) and headers[idx]
                                else f"column_{idx + 1}"
                            )
                            paired.append(f"{key}: {value}")
                    text = "; ".join(paired or values_for_text)
                    start_col = cells[0]["address"]
                    end_col = cells[-1]["address"]
                    locator = f"sheet:{sheet_index}:row:{row_index}"
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.TABLE,
                        locator,
                        text=text,
                        role="TABLE_ROW",
                        source_locator={
                            "sheetName": sheet_name,
                            "sheetIndex": sheet_index,
                            "rowIndex": row_index,
                            "rangeAddress": f"{start_col}:{end_col}",
                        },
                        metadata={
                            "cells": cells,
                            "hiddenSheet": ws_formula.sheet_state != "visible",
                            "headerValues": headers or [],
                        },
                    )
                if emitted_rows == 0:
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.OTHER,
                        f"sheet:{sheet_index}:empty",
                        role="EMPTY_SHEET",
                        source_locator={"sheetName": sheet_name, "sheetIndex": sheet_index},
                    )
            return _doc(source, context, "openpyxl-structured", elements, [], [])
        finally:
            formula_wb.close()
            value_wb.close()


class PptxDocumentHandler:
    format_name = "PPTX"

    def parse(self, source: AcquiredSource, context: ParseContext):
        presentation = Presentation(str(source.local_path))
        if len(presentation.slides) > context.limits.max_slides:
            raise ParserError("PARSER_SLIDE_LIMIT", "presentation slide limit exceeded")
        elements: list[DocumentElement] = []
        candidates: list[OcrCandidate] = []
        shape_count = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            ordered_shapes = sorted(
                list(slide.shapes),
                key=lambda shape: (
                    0
                    if getattr(shape, "is_placeholder", False)
                    and getattr(shape.placeholder_format, "type", None) == 1
                    else 1,
                    int(getattr(shape, "top", 0)),
                    int(getattr(shape, "left", 0)),
                    int(getattr(shape, "shape_id", 0)),
                ),
            )
            for shape_index, shape in enumerate(ordered_shapes):
                shape_count += 1
                if shape_count > context.limits.max_shapes:
                    raise ParserError("PARSER_SHAPE_LIMIT", "presentation shape limit exceeded")
                shape_id = int(getattr(shape, "shape_id", shape_index))
                locator = f"slide:{slide_number}:shape:{shape_id}"
                geometry = SourceGeometry(
                    page_number=slide_number,
                    x0=float(getattr(shape, "left", 0)),
                    y0=float(getattr(shape, "top", 0)),
                    x1=float(getattr(shape, "left", 0) + getattr(shape, "width", 0)),
                    y1=float(getattr(shape, "top", 0) + getattr(shape, "height", 0)),
                    page_width=float(presentation.slide_width),
                    page_height=float(presentation.slide_height),
                    coordinate_space="emu-top-left",
                )
                source_locator = {
                    "slideNumber": slide_number,
                    "shapeId": shape_id,
                    "shapeIndex": shape_index,
                }
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.TABLE,
                        locator,
                        text="\n".join(" | ".join(row) for row in rows),
                        geometry=geometry,
                        source_locator=source_locator,
                        metadata={
                            "rows": rows,
                            "rowCount": len(rows),
                            "columnCount": len(rows[0]) if rows else 0,
                        },
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    element = _append(
                        elements,
                        source,
                        context,
                        ElementType.IMAGE,
                        locator,
                        geometry=geometry,
                        source_locator=source_locator,
                        role="SLIDE_IMAGE",
                    )
                    candidates.append(
                        OcrCandidate(
                            candidate_id=f"ocr:{element.element_id}",
                            scope="EMBEDDED_IMAGE",
                            page_number=slide_number,
                            reason="pptx_picture",
                            element_id=element.element_id,
                            geometry=geometry,
                            metadata=source_locator,
                        )
                    )
                elif getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if not text:
                        continue
                    placeholder_type = None
                    if getattr(shape, "is_placeholder", False):
                        placeholder_type = str(shape.placeholder_format.type)
                    is_title = shape is slide.shapes.title
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.HEADING if is_title else ElementType.PARAGRAPH,
                        locator,
                        text=text,
                        geometry=geometry,
                        source_locator=source_locator,
                        role="SLIDE_TITLE" if is_title else "SLIDE_TEXT",
                        metadata={"placeholderType": placeholder_type},
                    )
            if slide.has_notes_slide:
                notes_text = (
                    "\n".join(
                        shape.text.strip()
                        for shape in slide.notes_slide.notes_text_frame.paragraphs
                        if getattr(shape, "text", "").strip()
                    )
                    if getattr(slide.notes_slide, "notes_text_frame", None)
                    else ""
                )
                if notes_text:
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.PARAGRAPH,
                        f"slide:{slide_number}:notes",
                        text=notes_text,
                        role="SPEAKER_NOTES",
                        source_locator={"slideNumber": slide_number, "notesScope": True},
                    )
        return _doc(source, context, "python-pptx-structured", elements, candidates, [])


class CsvDocumentHandler:
    format_name = "CSV"

    def parse(self, source: AcquiredSource, context: ParseContext):
        elements: list[DocumentElement] = []
        with source.local_path.open("r", encoding="utf-8", newline="") as stream:
            sample = stream.read(65_536)
            stream.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",;\t")
            except csv.Error as exc:
                raise ParserError("PARSER_CSV_DIALECT", "CSV dialect cannot be determined") from exc
            reader = csv.reader(stream, dialect)
            headers: list[str] | None = None
            for row_index, row in enumerate(reader, start=1):
                if row_index > context.limits.max_rows:
                    raise ParserError("PARSER_ROW_LIMIT", "CSV row limit exceeded")
                if len(row) > context.limits.max_columns:
                    raise ParserError("PARSER_COLUMN_LIMIT", "CSV column limit exceeded")
                if any(len(cell) > context.limits.max_cell_chars for cell in row):
                    raise ParserError(
                        "PARSER_CELL_TEXT_LIMIT", "CSV cell text exceeds configured limit"
                    )
                if headers is None:
                    headers = row
                paired = []
                for index, value in enumerate(row):
                    key = (
                        headers[index]
                        if headers and index < len(headers) and headers[index]
                        else f"column_{index + 1}"
                    )
                    paired.append(f"{key}: {value}")
                _append(
                    elements,
                    source,
                    context,
                    ElementType.TABLE,
                    f"row:{row_index}",
                    text="; ".join(paired),
                    role="TABLE_ROW",
                    source_locator={"rowIndex": row_index, "columnStart": 1, "columnEnd": len(row)},
                    metadata={
                        "cells": row,
                        "headers": headers or [],
                        "delimiter": dialect.delimiter,
                        "quoteChar": dialect.quotechar,
                    },
                )
        return _doc(source, context, "csv-stream-v1", elements, [], [])


def _dom_path(tag: Tag) -> str:
    parts: list[str] = []
    node: Tag | None = tag
    while node is not None and getattr(node, "name", None):
        parent = node.parent if isinstance(node.parent, Tag) else None
        index = 1
        if parent is not None:
            siblings = [s for s in parent.find_all(node.name, recursive=False)]
            if len(siblings) > 1:
                index = siblings.index(node) + 1
        parts.append(f"{node.name}[{index}]")
        node = parent
    return "/" + "/".join(reversed(parts))


class HtmlDocumentHandler:
    format_name = "HTML"

    def parse(self, source: AcquiredSource, context: ParseContext):
        markup = source.local_path.read_text(encoding="utf-8")
        soup = BeautifulSoup(markup, "lxml")
        for tag in soup(["script", "style", "noscript", "template"]):
            tag.decompose()
        nodes = soup.find_all(True)
        if len(nodes) > context.limits.max_dom_nodes:
            raise ParserError("PARSER_DOM_LIMIT", "HTML DOM node limit exceeded")
        elements: list[DocumentElement] = []
        for tag in nodes:
            name = tag.name.lower()
            if name not in {
                "title",
                "h1",
                "h2",
                "h3",
                "h4",
                "h5",
                "h6",
                "p",
                "li",
                "table",
                "pre",
                "code",
                "img",
            }:
                continue
            locator = _dom_path(tag)
            source_locator = {"domPath": locator, "elementId": tag.get("id")}
            if name == "img":
                _append(
                    elements,
                    source,
                    context,
                    ElementType.IMAGE,
                    locator,
                    role="HTML_IMAGE_REFERENCE",
                    source_locator=source_locator,
                    metadata={"src": tag.get("src"), "alt": tag.get("alt")},
                )
                continue
            text = tag.get_text(" ", strip=True)
            if not text:
                continue
            if name == "title" or name.startswith("h"):
                kind = ElementType.HEADING
                level = int(name[1]) if name.startswith("h") else 1
            elif name == "li":
                kind = ElementType.LIST_ITEM
                level = None
            elif name == "table":
                kind = ElementType.TABLE
                level = None
            elif name in {"pre", "code"}:
                kind = ElementType.CODE_BLOCK
                level = None
            else:
                kind = ElementType.PARAGRAPH
                level = None
            _append(
                elements,
                source,
                context,
                kind,
                locator,
                text=text,
                level=level,
                source_locator=source_locator,
            )
        return _doc(source, context, "beautifulsoup-lxml", elements, [], [])


class RtfDocumentHandler:
    format_name = "RTF"

    def parse(self, source: AcquiredSource, context: ParseContext):
        raw = source.local_path.read_text(encoding="latin-1")
        text = rtf_to_text(raw)
        elements: list[DocumentElement] = []
        for paragraph_index, paragraph in enumerate(
            p.strip() for p in text.splitlines() if p.strip()
        ):
            _append(
                elements,
                source,
                context,
                ElementType.PARAGRAPH,
                f"paragraph:{paragraph_index}",
                text=paragraph,
                source_locator={"paragraphIndex": paragraph_index},
            )
        return _doc(source, context, "striprtf-v1", elements, [], ["RTF_STRUCTURE_PARTIAL"])


class OdtDocumentHandler:
    format_name = "ODT"

    def parse(self, source: AcquiredSource, context: ParseContext):
        document = load_odf(str(source.local_path))
        elements: list[DocumentElement] = []
        index = 0

        def walk(node: Any, section_path: tuple[str, ...] = ()) -> None:
            nonlocal index
            qname = getattr(node, "qname", None)
            local = qname[1] if qname else None
            if local in {"h", "p", "list", "table"}:
                text = teletype.extractText(node).strip()
                if text:
                    if local == "h":
                        kind = ElementType.HEADING
                    elif local == "list":
                        kind = ElementType.LIST
                    elif local == "table":
                        kind = ElementType.TABLE
                    else:
                        kind = ElementType.PARAGRAPH
                    _append(
                        elements,
                        source,
                        context,
                        kind,
                        f"node:{index}",
                        text=text,
                        section_path=section_path,
                        source_locator={"odfNodeIndex": index, "odfType": local},
                    )
                    index += 1
                    return
            for child in getattr(node, "childNodes", ()):
                walk(child, section_path)

        walk(document.text)
        return _doc(source, context, "odfpy-structured", elements, [], [])


class EpubDocumentHandler:
    format_name = "EPUB"

    def parse(self, source: AcquiredSource, context: ParseContext):
        book = epub.read_epub(str(source.local_path))
        if len(book.spine) > context.limits.max_epub_spine_items:
            raise ParserError("PARSER_EPUB_SPINE_LIMIT", "EPUB spine item limit exceeded")
        elements: list[DocumentElement] = []
        item_by_id = {item.id: item for item in book.get_items()}
        for spine_index, spine_entry in enumerate(book.spine):
            item_id = spine_entry[0]
            item = item_by_id.get(item_id)
            if item is None or item.get_type() != ITEM_DOCUMENT:
                continue
            soup = BeautifulSoup(item.get_content(), "lxml")
            for tag in soup(["script", "style", "noscript", "template"]):
                tag.decompose()
            for local_index, tag in enumerate(
                soup.find_all(
                    ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table", "pre", "img"]
                )
            ):
                locator = f"spine:{spine_index}:node:{local_index}"
                source_locator = {
                    "spineIndex": spine_index,
                    "href": item.file_name,
                    "domPath": _dom_path(tag),
                }
                if tag.name == "img":
                    _append(
                        elements,
                        source,
                        context,
                        ElementType.IMAGE,
                        locator,
                        role="EPUB_IMAGE_REFERENCE",
                        source_locator=source_locator,
                        metadata={"src": tag.get("src"), "alt": tag.get("alt")},
                    )
                    continue
                text = tag.get_text(" ", strip=True)
                if not text:
                    continue
                if tag.name.startswith("h"):
                    kind = ElementType.HEADING
                    level = int(tag.name[1])
                elif tag.name == "li":
                    kind = ElementType.LIST_ITEM
                    level = None
                elif tag.name == "table":
                    kind = ElementType.TABLE
                    level = None
                elif tag.name == "pre":
                    kind = ElementType.CODE_BLOCK
                    level = None
                else:
                    kind = ElementType.PARAGRAPH
                    level = None
                _append(
                    elements,
                    source,
                    context,
                    kind,
                    locator,
                    text=text,
                    level=level,
                    source_locator=source_locator,
                )
        return _doc(source, context, "ebooklib-html-structured", elements, [], [])
