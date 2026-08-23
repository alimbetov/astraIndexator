from __future__ import annotations

from bs4 import BeautifulSoup
from ebooklib import ITEM_DOCUMENT, epub
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from astra_indexator.acquisition import AcquiredSource

from .base import ParseContext, ParserError
from .extended_handlers import _append, _dom_path
from .handlers import _doc
from .model import ElementType, OcrCandidate, SourceGeometry


class PptxDocumentHandlerV1:
    format_name = "PPTX"

    def parse(self, source: AcquiredSource, context: ParseContext):
        presentation = Presentation(str(source.local_path))
        if len(presentation.slides) > context.limits.max_slides:
            raise ParserError("PARSER_SLIDE_LIMIT", "presentation slide limit exceeded")
        elements = []
        candidates: list[OcrCandidate] = []
        shape_count = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape_id = slide.shapes.title.shape_id if slide.shapes.title is not None else None
            ordered_shapes = sorted(
                list(slide.shapes),
                key=lambda shape: (
                    0 if getattr(shape, "shape_id", None) == title_shape_id else 1,
                    int(getattr(shape, "top", 0)),
                    int(getattr(shape, "left", 0)),
                    int(getattr(shape, "shape_id", 0)),
                ),
            )
            for shape_index, shape in enumerate(ordered_shapes):
                shape_count += 1
                if shape_count > context.limits.max_shapes:
                    raise ParserError("PARSER_SHAPE_LIMIT", "presentation shape limit exceeded")
                shape_id = int(getattr(shape, "shape_id", shape_index))
                locator = f"slide:{slide_number}:shape:{shape_id}"
                geometry = SourceGeometry(
                    page_number=slide_number,
                    x0=float(getattr(shape, "left", 0)),
                    y0=float(getattr(shape, "top", 0)),
                    x1=float(getattr(shape, "left", 0) + getattr(shape, "width", 0)),
                    y1=float(getattr(shape, "top", 0) + getattr(shape, "height", 0)),
                    page_width=float(presentation.slide_width),
                    page_height=float(presentation.slide_height),
                    coordinate_space="emu-top-left",
                )
                source_locator = {"slideNumber": slide_number, "shapeId": shape_id, "shapeIndex": shape_index}
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    _append(
                        elements, source, context, ElementType.TABLE, locator,
                        text="\n".join(" | ".join(row) for row in rows),
                        geometry=geometry, source_locator=source_locator,
                        metadata={"rows": rows, "rowCount": len(rows), "columnCount": len(rows[0]) if rows else 0},
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    element = _append(
                        elements, source, context, ElementType.IMAGE, locator,
                        geometry=geometry, source_locator=source_locator, role="SLIDE_IMAGE",
                    )
                    candidates.append(
                        OcrCandidate(
                            candidate_id=f"ocr:{element.element_id}",
                            scope="EMBEDDED_IMAGE",
                            page_number=slide_number,
                            reason="pptx_picture",
                            element_id=element.element_id,
                            geometry=geometry,
                            metadata=source_locator,
                        )
                    )
                elif getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if not text:
                        continue
                    placeholder_type = str(shape.placeholder_format.type) if getattr(shape, "is_placeholder", False) else None
                    is_title = shape_id == title_shape_id
                    _append(
                        elements, source, context,
                        ElementType.HEADING if is_title else ElementType.PARAGRAPH,
                        locator,
                        text=text,
                        geometry=geometry,
                        source_locator=source_locator,
                        role="SLIDE_TITLE" if is_title else "SLIDE_TEXT",
                        metadata={"placeholderType": placeholder_type},
                    )
            notes_slide = slide.notes_slide if slide.has_notes_slide else None
            notes_frame = getattr(notes_slide, "notes_text_frame", None) if notes_slide is not None else None
            if notes_frame is not None:
                notes_text = "\n".join(p.text.strip() for p in notes_frame.paragraphs if p.text.strip())
                if notes_text:
                    _append(
                        elements, source, context, ElementType.PARAGRAPH,
                        f"slide:{slide_number}:notes", text=notes_text,
                        role="SPEAKER_NOTES",
                        source_locator={"slideNumber": slide_number, "notesScope": True},
                    )
        return _doc(source, context, "python-pptx-structured", elements, candidates, [])


class EpubDocumentHandlerV1:
    format_name = "EPUB"

    def parse(self, source: AcquiredSource, context: ParseContext):
        book = epub.read_epub(str(source.local_path))
        if len(book.spine) > context.limits.max_epub_spine_items:
            raise ParserError("PARSER_EPUB_SPINE_LIMIT", "EPUB spine item limit exceeded")
        elements = []
        item_by_id = {item.id: item for item in book.get_items()}
        for spine_index, spine_entry in enumerate(book.spine):
            item_id = spine_entry[0]
            item = item_by_id.get(item_id)
            if item is None or item.get_type() != ITEM_DOCUMENT:
                continue
            href = item.file_name or ""
            if item_id.lower() == "nav" or href.lower().endswith(("nav.xhtml", "nav.html")):
                continue
            soup = BeautifulSoup(item.get_content(), "lxml")
            for tag in soup(["script", "style", "noscript", "template"]):
                tag.decompose()
            for local_index, tag in enumerate(
                soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table", "pre", "img"])
            ):
                locator = f"spine:{spine_index}:node:{local_index}"
                source_locator = {"spineIndex": spine_index, "href": href, "domPath": _dom_path(tag)}
                if tag.name == "img":
                    _append(
                        elements, source, context, ElementType.IMAGE, locator,
                        role="EPUB_IMAGE_REFERENCE", source_locator=source_locator,
                        metadata={"src": tag.get("src"), "alt": tag.get("alt")},
                    )
                    continue
                text = tag.get_text(" ", strip=True)
                if not text:
                    continue
                if tag.name.startswith("h"):
                    kind = ElementType.HEADING
                    level = int(tag.name[1])
                elif tag.name == "li":
                    kind = ElementType.LIST_ITEM
                    level = None
                elif tag.name == "table":
                    kind = ElementType.TABLE
                    level = None
                elif tag.name == "pre":
                    kind = ElementType.CODE_BLOCK
                    level = None
                else:
                    kind = ElementType.PARAGRAPH
                    level = None
                _append(elements, source, context, kind, locator, text=text, level=level, source_locator=source_locator)
        return _doc(source, context, "ebooklib-html-structured", elements, [], [])
