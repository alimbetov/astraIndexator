from __future__ import annotations

import hashlib
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

from ebooklib import epub
from odf.opendocument import OpenDocumentText
from odf.text import H, P
from openpyxl import Workbook
from pptx import Presentation

from astra_indexator.acquisition import AcquiredSource, SafeAcquisitionService
from astra_indexator.parser import (
    DocumentParserService,
    ElementType,
    ParseContext,
    default_registry,
)
from astra_indexator.storage import StorageRef


def _source(path: Path, fmt: str) -> AcquiredSource:
    payload = path.read_bytes()
    return AcquiredSource(
        source_ref=StorageRef.parse(f"seaweed://documents/{path.name}"),
        local_path=path,
        original_file_name=path.name,
        detected_format=fmt,
        detected_content_type="application/octet-stream",
        size_bytes=len(payload),
        sha256=hashlib.sha256(payload).hexdigest(),
        etag=None,
        version_id=None,
        validation_profile="default-v1",
        warnings=(),
        acquired_at=datetime.now(timezone.utc),
    )


def _parse(path: Path, fmt: str):
    source = _source(path, fmt)
    context = ParseContext(uuid4(), uuid4(), uuid4(), 1, source.sha256)
    return DocumentParserService(default_registry()).parse(source, context)


def _validator(tmp_path: Path) -> SafeAcquisitionService:
    return SafeAcquisitionService(storage=object(), workspace_root=tmp_path)  # type: ignore[arg-type]


def test_xlsx_preserves_sheet_row_cell_and_formula_provenance(tmp_path: Path) -> None:
    path = tmp_path / "report.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Month", "Amount", "Formula"])
    ws.append(["Jan", 100, "=B2*2"])
    wb.save(path)

    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "XLSX"
    result = _parse(path, fmt)
    rows = [element for element in result.elements if element.role == "TABLE_ROW"]
    assert rows
    assert rows[1].source_locator["sheetName"] == "Revenue"
    assert rows[1].source_locator["rowIndex"] == 2
    assert any(cell["formulaPresent"] for cell in rows[1].metadata["cells"])


def test_pptx_preserves_slide_shape_and_table_provenance(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Architecture"
    table = slide.shapes.add_table(2, 2, 1000000, 2000000, 4000000, 1500000).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    prs.save(path)

    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "PPTX"
    result = _parse(path, fmt)
    assert any(e.type == ElementType.HEADING and e.text == "Architecture" for e in result.elements)
    table_element = next(e for e in result.elements if e.type == ElementType.TABLE)
    assert table_element.source_locator["slideNumber"] == 1
    assert "shapeId" in table_element.source_locator


def test_csv_is_streamed_as_header_associated_structured_rows(tmp_path: Path) -> None:
    path = tmp_path / "people.csv"
    path.write_text("name;age;city\nAruzhan;31;Almaty\nRuslan;44;Almaty\n", encoding="utf-8")
    fmt, _, warnings = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "CSV"
    assert any(w.startswith("CSV_DELIMITER") for w in warnings)
    result = _parse(path, fmt)
    assert all(e.type == ElementType.TABLE for e in result.elements)
    assert result.elements[1].source_locator["rowIndex"] == 2
    assert "name: Aruzhan" in (result.elements[1].text or "")


def test_html_is_offline_structured_dom_without_script_content(tmp_path: Path) -> None:
    path = tmp_path / "page.html"
    path.write_text(
        "<!doctype html><html><body><h1>Policy</h1><p id='p1'>Text</p>"
        "<script>fetch('https://example.com')</script><img src='https://example.com/a.png' alt='A'></body></html>",
        encoding="utf-8",
    )
    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "HTML"
    result = _parse(path, fmt)
    assert any(e.type == ElementType.HEADING and e.text == "Policy" for e in result.elements)
    assert all("fetch(" not in (e.text or "") for e in result.elements)
    image = next(e for e in result.elements if e.type == ElementType.IMAGE)
    assert image.metadata["src"].startswith("https://")
    assert not result.ocr_candidates  # remote resources are recorded, never fetched


def test_rtf_admission_and_paragraph_conversion(tmp_path: Path) -> None:
    path = tmp_path / "sample.rtf"
    path.write_bytes(rb"{\rtf1\ansi First paragraph.\par Second paragraph.}")
    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "RTF"
    result = _parse(path, fmt)
    assert result.elements
    assert all(e.type == ElementType.PARAGRAPH for e in result.elements)
    assert "RTF_STRUCTURE_PARTIAL" in result.quality.warnings


def test_odt_preserves_logical_structure_without_fake_pages(tmp_path: Path) -> None:
    path = tmp_path / "sample.odt"
    doc = OpenDocumentText()
    doc.text.addElement(H(outlinelevel=1, text="Guide"))
    doc.text.addElement(P(text="Kazakh and Russian content"))
    doc.save(str(path), addsuffix=False)
    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "ODT"
    result = _parse(path, fmt)
    assert any(e.type == ElementType.HEADING and e.text == "Guide" for e in result.elements)
    assert all(e.geometry is None for e in result.elements)


def test_epub_uses_spine_order_and_href_provenance(tmp_path: Path) -> None:
    path = tmp_path / "book.epub"
    book = epub.EpubBook()
    book.set_identifier("book-1")
    book.set_title("Test")
    book.set_language("en")
    c1 = epub.EpubHtml(title="One", file_name="one.xhtml", lang="en")
    c1.content = "<h1>One</h1><p>First chapter.</p>"
    c2 = epub.EpubHtml(title="Two", file_name="two.xhtml", lang="en")
    c2.content = "<h1>Two</h1><p>Second chapter.</p>"
    book.add_item(c1)
    book.add_item(c2)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", c1, c2]
    epub.write_epub(str(path), book)

    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "EPUB"
    result = _parse(path, fmt)
    headings = [e for e in result.elements if e.type == ElementType.HEADING]
    assert [e.text for e in headings] == ["One", "Two"]
    assert headings[0].source_locator["spineIndex"] < headings[1].source_locator["spineIndex"]
    assert headings[0].source_locator["href"] == "one.xhtml"


def test_plain_delimited_text_is_not_promoted_to_csv_without_stable_evidence(
    tmp_path: Path,
) -> None:
    path = tmp_path / "note.txt"
    path.write_text("alpha,beta\nonly one narrative line afterwards\n", encoding="utf-8")
    fmt, _, _ = _validator(tmp_path)._validate(path, path.name)
    assert fmt == "TXT"
